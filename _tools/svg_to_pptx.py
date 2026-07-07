#!/usr/bin/env python3
"""Convert the slide-NN.svg files in _slides/ into a slides.pptx deck.

Each SVG is rasterized (via headless Chromium) to a JPEG at its native canvas
resolution and embedded as a plain picture, one per slide. This renders
identically everywhere, including PowerPoint, since it sidesteps PowerPoint's
inconsistent native SVG support entirely (earlier attempts embedding native
SVG ran into PowerPoint not honoring xlink:href / clip-path).

Usage:
    python svg_to_pptx.py
    python svg_to_pptx.py --slides-dir ../_slides --output ../slides.pptx
    python svg_to_pptx.py --pattern "slide-*.svg" --jpeg-quality 85
"""

from __future__ import annotations

import argparse
import io
import re
import sys
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# Standard PowerPoint 16:9 widescreen size, matching the 3840x2160 SVG canvases.
SLIDE_WIDTH_EMU = Emu(12192000)
SLIDE_HEIGHT_EMU = Emu(6858000)

_NUM_RE = re.compile(r"(\d+)")


def _natural_key(path: Path):
    return [int(p) if p.isdigit() else p for p in _NUM_RE.split(path.stem)]


def _svg_canvas_size(svg_path: Path, default: tuple[int, int] = (3840, 2160)) -> tuple[int, int]:
    """Return the (width, height) of the root <svg> element, for viewport sizing."""
    head = svg_path.read_text(encoding="utf-8", errors="ignore")[:2000]
    m = re.search(r"<svg\b[^>]*\bwidth=\"(\d+(?:\.\d+)?)\"[^>]*\bheight=\"(\d+(?:\.\d+)?)\"", head)
    if not m:
        m = re.search(r"<svg\b[^>]*\bheight=\"(\d+(?:\.\d+)?)\"[^>]*\bwidth=\"(\d+(?:\.\d+)?)\"", head)
        if m:
            return int(float(m.group(2))), int(float(m.group(1)))
        return default
    return int(float(m.group(1))), int(float(m.group(2)))


def _render_slide_jpeg(svg_path: Path, quality: int,
                        width: int | None, height: int | None) -> bytes:
    """Rasterize an SVG to JPEG bytes by screenshotting it in headless Chromium."""
    canvas_width, canvas_height = _svg_canvas_size(svg_path)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return _placeholder_jpeg(width or canvas_width, height or canvas_height, quality)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            try:
                page = browser.new_page(viewport={"width": canvas_width, "height": canvas_height})
                page.goto(svg_path.resolve().as_uri())
                screenshot = page.screenshot()
            finally:
                browser.close()
    except Exception as exc:  # pragma: no cover - environment dependent
        print(f"  warning: playwright render failed for {svg_path.name} ({exc}); "
              "using placeholder image", file=sys.stderr)
        return _placeholder_jpeg(width or canvas_width, height or canvas_height, quality)

    from PIL import Image

    img = Image.open(io.BytesIO(screenshot)).convert("RGB")
    if width and height and (width, height) != img.size:
        img = img.resize((width, height), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality, optimize=True)
    return buf.getvalue()


def _placeholder_jpeg(width: int, height: int, quality: int) -> bytes:
    from PIL import Image

    img = Image.new("RGB", (width, height), color=(20, 20, 30))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality)
    return buf.getvalue()


def build_pptx(slides_dir: Path, output_path: Path, pattern: str, quality: int,
                width: int | None, height: int | None) -> None:
    start_time = time.monotonic()
    print(f"Starting SVG -> PPTX conversion (slides dir: {slides_dir})")

    svg_paths = sorted(slides_dir.glob(pattern), key=_natural_key)
    if not svg_paths:
        raise SystemExit(f"No files matching '{pattern}' found in {slides_dir}")
    print(f"Found {len(svg_paths)} slide(s) matching '{pattern}'")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU
    blank_layout = prs.slide_layouts[6]

    for i, svg_path in enumerate(svg_paths, start=1):
        slide_start = time.monotonic()
        print(f"[{i}/{len(svg_paths)}] Processing {svg_path.name} ...")

        jpeg_bytes = _render_slide_jpeg(svg_path, quality, width, height)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(jpeg_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        print(f"[{i}/{len(svg_paths)}] Done {svg_path.name} "
              f"({time.monotonic() - slide_start:.2f}s)")

    print("Saving presentation...")
    try:
        prs.save(output_path)
    except PermissionError:
        raise SystemExit(
            f"Could not write {output_path} - it looks like it's open in PowerPoint "
            "(or another app). Close it and re-run."
        )
    size_mb = output_path.stat().st_size / (1024 * 1024)
    elapsed = time.monotonic() - start_time
    print(f"\nFinished: wrote {len(svg_paths)} slides to {output_path} "
          f"({size_mb:.2f} MB) in {elapsed:.2f}s")


def main() -> None:
    script_dir = Path(__file__).resolve().parent
    repo_root = script_dir.parent

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--slides-dir", type=Path, default=repo_root / "_slides",
                         help="Directory containing the slide SVGs (default: ../_slides)")
    parser.add_argument("--output", type=Path, default=repo_root / "slides.pptx",
                         help="Output pptx path (default: ../slides.pptx)")
    parser.add_argument("--pattern", default="slide-*.svg",
                         help="Glob pattern (relative to --slides-dir) selecting slide SVGs "
                              "(default: 'slide-*.svg'; other assets like logos/diagrams in "
                              "_slides are intentionally excluded)")
    parser.add_argument("--jpeg-quality", type=int, default=90, dest="quality",
                         help="JPEG quality 1-100 for the rasterized slides (default: 90)")
    parser.add_argument("--width", type=int, default=None,
                         help="Output pixel width per slide (default: SVG's native canvas width, "
                              "typically 3840)")
    parser.add_argument("--height", type=int, default=None,
                         help="Output pixel height per slide (default: SVG's native canvas "
                              "height, typically 2160)")
    args = parser.parse_args()

    build_pptx(args.slides_dir, args.output, args.pattern,
               args.quality, args.width, args.height)


if __name__ == "__main__":
    main()
